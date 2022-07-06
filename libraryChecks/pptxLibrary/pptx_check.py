import pdb
import random
from pptx import Presentation
from libraryChecks.pptxLibrary.constants import *
from pptx.util import Inches


def add_shakespeare_quote_text():
    prs = Presentation(INPUT_FILEPATH)
    quote_slide = prs.slides.get(SHAKESPEAR_SLIDE_ID)
    quote_shape = quote_slide.placeholders[SHAKESPEAR_QUOTE_SHAPE_ID]
    quote_shape.text = "Cowards die many times before their deaths; The valiant never taste of death but once."
    prs.save(OUTPUT_FILEPATH)


def add_gandhi_quote_img():
    prs = Presentation(INPUT_FILEPATH)
    quote_slide = prs.slides.get(GANDHI_SLIDE_ID)
    pic = quote_slide.shapes.add_picture(GANDHI_QUOTE_IMAGE, Inches(2.5), Inches(2.4), width=Inches(9), height=Inches(4.8))
    prs.save(OUTPUT_FILEPATH)


def add_quotes_table():
    prs = Presentation(INPUT_FILEPATH)
    table_slide = prs.slides.get(TABLE_SLIDE_ID)
    table_gf = table_slide.shapes.add_table(rows=3, cols=4, left=Inches(2.5), top=Inches(2.4), width=Inches(9), height=Inches(4.8))
    table = table_gf.table
    pdb.set_trace()
    for cell in table.iter_cells():
        cell.text = str(random.randrange(100))
    prs.save(OUTPUT_FILEPATH)

add_quotes_table()