import pytest
from pptx import Presentation
from libraryChecks.pptxLibrary.constants import *
import pdb


def test_shakespeare_quote_text():
    prs = Presentation(OUTPUT_FILEPATH)
    text_list = []
    quote_slide = prs.slides.get(SHAKESPEAR_SLIDE_ID)
    pdb.set_trace()
    quotes_shape = quote_slide.shapes[SHAKESPEAR_QUOTE_SHAPE_ID]
    for paragraph in quotes_shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text_list.append(run.text)
    assert "".join(text_list) == SHAKESPEAR_QUOTE


def test_gandhi_quote_img():
    prs = Presentation(OUTPUT_FILEPATH)
    quote_slide = prs.slides.get(GANDHI_SLIDE_ID)
    quote_shape = quote_slide.shapes[GANDHI_QUOTE_SHAPE_ID]
    image_filename = '/Users/isibrahi/Personal/projects/BasicProgramming/libraryChecks/pptxLibrary/extracted_image.%s' % quote_shape.image.ext
    image_bytes = quote_shape.image.blob
    assert image_bytes


def test_quote_table():
    row_len =0
    col_len =0
    prs = Presentation(OUTPUT_FILEPATH)
    table_slide = prs.slides.get(TABLE_SLIDE_ID)
    for shape in table_slide.shapes:
        if shape.has_table:
            col_len = len(shape.table.columns)
            row_len = len(shape.table.rows)

    assert row_len == 3
    assert col_len == 4

test_quote_table()